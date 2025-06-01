from pptx import Presentation
import difflib

from pptx.util import Pt
# Target keys you want to extract values for
elems = {"Contractor": None, 'Consultant': None, 'LOA Date': None}

# Helper: find closest match from known keys


def find_closest_key(cell_text, threshold=0.7):
    matches = difflib.get_close_matches(
        cell_text, elems.keys(), n=1, cutoff=threshold)
    return matches[0] if matches else None

# Recursive shape scanner


def find_and_extract_text(shapes):
    for shape in shapes:
        if shape.shape_type == 6:  # GroupShape
            find_and_extract_text(shape.shapes)
        elif shape.has_table:
            for row in shape.table.rows:
                for cell_idx in range(len(row.cells) - 1):
                    key_cell = row.cells[cell_idx]
                    match = find_closest_key(key_cell.text.strip())
                    if match:
                        elems[match] = row.cells[cell_idx + 1]


# Load the presentation
prs = Presentation("Sample_Project.pptx")

# Access the second slide
slide = prs.slides[1]

# Perform extraction
find_and_extract_text(slide.shapes)
print(Pt(15))
# Output results
for k, v in elems.items():
    print(k)
    print(v.text_frame.paragraphs[0])
    for paragraph in v.text_frame.paragraphs:
        print(paragraph.runs)
        # for run in paragraph.runs:
        #
        #     print(run.font.size)  # Set font size to 18 pt
