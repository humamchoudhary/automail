import os
import json
import base64
import re
import io
from datetime import datetime
from flask import Flask, render_template, redirect, url_for, session, request, flash, send_file
import secrets

import glob
import google.generativeai as genai
# Gmail API imports
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import Flow
from googleapiclient.discovery import build

# PowerPoint generation imports
from pptx import Presentation
from pptx.util import Pt

import PyPDF2
from docx import Document
import tempfile

import difflib
app = Flask(__name__)
app.secret_key = secrets.token_hex(16)

# OAuth 2.0 configuration
SCOPES = ['https://www.googleapis.com/auth/gmail.modify']
CLIENT_SECRETS_FILE = "credentials.json"


class GmailPPTProcessor:
    def __init__(self):
        self.processed_emails_file = 'processed_emails.json'
        self.config_file = 'config.json'
        self.processed_emails = self.load_processed_emails()
        self.config = self.load_config()

        # Create necessary directories
        os.makedirs('processed', exist_ok=True)
        os.makedirs('templates', exist_ok=True)
        os.makedirs('static', exist_ok=True)

        # Configure Gemini
        api_key = os.environ.get('GEMINI_API_KEY')
        print(api_key)
        genai.configure(api_key=api_key)
        self.model = genai.GenerativeModel('gemma-3-4b-it')

    def load_config(self):
        """Load configuration from JSON file"""
        default_config = {
            "search_text": "GENERATE PPT",
            "ppt_template": "template.pptx",
            "title_slide_index": 0,
            "content_slide_index": 1
        }

        try:
            with open(self.config_file, 'r') as f:
                config = json.load(f)
                for key, value in default_config.items():
                    if key not in config:
                        config[key] = value
                return config
        except FileNotFoundError:
            with open(self.config_file, 'w') as f:
                json.dump(default_config, f, indent=4)
            return default_config

    def load_processed_emails(self):
        """Load list of already processed email IDs"""
        try:
            with open(self.processed_emails_file, 'r') as f:
                return set(json.load(f))
        except FileNotFoundError:
            return set()

    def save_processed_emails(self):
        """Save list of processed email IDs"""
        with open(self.processed_emails_file, 'w') as f:
            json.dump(list(self.processed_emails), f)

    def get_gmail_service(self, credentials):
        """Create Gmail service with provided credentials"""
        return build('gmail', 'v1', credentials=credentials)

    def get_unread_emails(self, service):
        """Get all unread emails from Gmail"""
        try:
            results = service.users().messages().list(
                userId='me', q='is:unread'
            ).execute()

            messages = results.get('messages', [])
            return messages

        except Exception as error:
            print(f"Error fetching emails: {error}")
            return []

    def get_email_content(self, service, message_id):
        """Get the content of a specific email"""
        try:
            message = service.users().messages().get(
                userId='me', id=message_id, format='full'
            ).execute()

            headers = message['payload'].get('headers', [])
            subject = next(
                (h['value'] for h in headers if h['name'] == 'Subject'), 'No Subject')
            sender = next(
                (h['value'] for h in headers if h['name'] == 'From'), 'Unknown Sender')

            body = self.extract_email_body(message['payload'])
            attachments = self.extract_attachments(
                service, message['payload'], message_id)

            return {
                'id': message_id,
                'subject': subject,
                'sender': sender,
                'body': body,
                'attachments': attachments
            }

        except Exception as error:
            print(f"Error getting email content: {error}")
            return None

    def extract_email_body(self, payload):
        """Extract text content from email payload"""
        body = ""

        if 'parts' in payload:
            for part in payload['parts']:
                if part['mimeType'] == 'text/plain':
                    data = part['body']['data']
                    body += base64.urlsafe_b64decode(data).decode('utf-8')
                elif part['mimeType'] == 'text/html':
                    data = part['body']['data']
                    html_body = base64.urlsafe_b64decode(data).decode('utf-8')
                    body += re.sub('<[^<]+?>', '', html_body)
        else:
            if payload['mimeType'] == 'text/plain':
                data = payload['body']['data']
                body = base64.urlsafe_b64decode(data).decode('utf-8')

        return body

    def extract_attachments(self, service, payload, message_id):
        """Extract text from email attachments (PDF, DOCX, PPTX, TXT)"""
        attachments_text = []

        def process_parts(parts):
            for part in parts:
                if 'filename' in part and part['filename']:
                    if part['body'].get('attachmentId'):
                        try:
                            attachment = service.users().messages().attachments().get(
                                userId='me', messageId=message_id, id=part['body']['attachmentId']
                            ).execute()

                            file_data = base64.urlsafe_b64decode(
                                attachment['data'])
                            filename = part['filename'].lower()

                            # Create a temporary file to process
                            with tempfile.NamedTemporaryFile(delete=False) as temp_file:
                                temp_file.write(file_data)
                                temp_file_path = temp_file.name

                            try:
                                if filename.endswith('.pdf'):
                                    text = self.extract_text_from_pdf(
                                        temp_file_path)
                                elif filename.endswith('.docx'):
                                    text = self.extract_text_from_docx(
                                        temp_file_path)
                                elif filename.endswith('.pptx'):
                                    text = self.extract_text_from_pptx(
                                        temp_file_path)
                                elif filename.endswith('.txt'):
                                    text = self.extract_text_from_txt(
                                        temp_file_path)
                                else:
                                    text = f"Unsupported file type: {filename}"

                                if text:
                                    attachments_text.append(
                                        f"Attachment: {filename}\n{text}")
                            finally:
                                # Clean up the temporary file
                                try:
                                    os.unlink(temp_file_path)
                                except:
                                    pass

                        except Exception as e:
                            print(f"Error processing attachment {
                                  part['filename']}: {e}")

                if 'parts' in part:
                    process_parts(part['parts'])

        if 'parts' in payload:
            process_parts(payload['parts'])

        return attachments_text

    def extract_text_from_pdf(self, file_path):
        """Extract text from PDF file"""
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return text
        except Exception as e:
            print(f"Error extracting text from PDF: {e}")
            return None

    def extract_text_from_docx(self, file_path):
        """Extract text from DOCX file"""
        try:
            doc = Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])
        except Exception as e:
            print(f"Error extracting text from DOCX: {e}")
            return None

    def extract_text_from_pptx(self, file_path):
        """Extract text from PPTX file"""
        try:
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            print(f"Error extracting text from PPTX: {e}")
            return None

    def extract_text_from_txt(self, file_path):
        """Extract text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            print(f"Error extracting text from TXT: {e}")
            return None

    def check_for_trigger_text(self, email_content):
        """Check if email contains contractor/consultant trigger keywords"""
        trigger_keywords = ['contractor', 'consultant', 'loa date', 'letter of authorization',
                            'vendor', 'service provider', 'contract', 'agreement']

        # Check email body
        content_to_check = email_content['body'].lower()

        # Check attachments
        for attachment_text in email_content['attachments']:
            content_to_check += "\n" + attachment_text.lower()

        for keyword in trigger_keywords:
            if keyword in content_to_check:
                return True

        return False

    def process_with_gemini(self, email_content):
        """Process email content with Gemini to extract contractor/consultant information"""
        try:
            # Prepare the content for Gemini
            full_content = f"""
            Email Subject: {email_content['subject']}
            Email Body: {email_content['body']}
            """

            if email_content['attachments']:
                full_content += f"\nAttachments Content: {
                    ' '.join(email_content['attachments'])}"

            # Create prompt for Gemini
            prompt = f"""
            Analyze the following email content and extract information related to contractors, consultants, and LOA (Letter of Authorization) dates.

            Please extract and return the following information if available:
            1. Contractor names and details
            2. Consultant names and details
            3. LOA dates mentioned
            4. Project or work descriptions
            5. Any other relevant contract or authorization information

            Email Content:
            {full_content}

            Please format your response as structured data with clear headings for each type of information found.
            """

            # Generate response using Gemini
            response = self.model.generate_content(prompt)
            return response.text

        except Exception as e:
            print(f"Error processing with Gemini: {e}")
            return "Error processing content with Gemini"

    def extract_info_with_gemini(self, text_content):
        """Use Gemini to extract structured information from text content"""
        try:
            prompt = f"""
            Analyze the following content and extract specific information in JSON format.
            Look for the following details regardless of how they are phrased in the text:

            1. Contractor information (could be called vendor, supplier, executing company, etc.)
            2. Consultant information (could be called advisor, supervisor, designer, etc.)
            3. LOA Date (could be called Letter of Authorization date, approval date, effective date, etc.)

            Return the information in this exact JSON format:
            {{
                "contractor": "name or details of the contractor",
                "consultant": "name or details of the consultant",
                "loa_date": "date in YYYY-MM-DD format if available"
            }}

            If any information is not found, use null for that field.

            Content to analyze:
            {text_content}
            """

            response = self.model.generate_content(prompt)

            try:
                # Extract JSON from Gemini's response
                json_str = response.text.replace(
                    '```json', '').replace('```', '').strip()
                extracted_data = json.loads(json_str)
                return extracted_data
            except json.JSONDecodeError:
                # Fallback if JSON parsing fails
                return {
                    "contractor": None,
                    "consultant": None,
                    "loa_date": None
                }

        except Exception as e:
            print(f"Error extracting info with Gemini: {e}")
            return {
                "contractor": None,
                "consultant": None,
                "loa_date": None
            }

    def extract_text_from_pptx(self, file_path):
        """Extract all text from PPTX file for Gemini analysis"""
        try:
            prs = Presentation(file_path)
            text_content = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                text_content.append(cell.text)

            return "\n".join(text_content)
        except Exception as e:
            print(f"Error extracting text from PPTX: {e}")
            return ""

    def extract_ppt_content(self, email_content):
        """Extract content for PPT generation using Gemini for info extraction"""
        # Get all text content from email and attachments
        full_text = f"EMAIL SUBJECT: {email_content['subject']}\n"
        full_text += f"EMAIL BODY: {email_content['body']}\n"

        if email_content['attachments']:
            full_text += "\nATTACHMENTS:\n" + \
                "\n".join(email_content['attachments'])

        # Extract information using Gemini
        extracted_info = self.extract_info_with_gemini(full_text)

        # Get detailed Gemini analysis
        gemini_analysis = self.process_with_gemini(email_content)

        content = {
            'title': email_content['subject'],
            'body': email_content['body'],
            'sender': email_content['sender'],
            'timestamp': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            'gemini_analysis': gemini_analysis,
            'extracted_info': extracted_info
        }

        if email_content['attachments']:
            content['attachments'] = '\n\n'.join(email_content['attachments'])

        return content

    def find_and_extract_text(self, shapes):
        def find_closest_key(cell_text, threshold=0.7):
            matches = difflib.get_close_matches(
                cell_text, self.elems.keys(), n=1, cutoff=threshold)
            return matches[0] if matches else None

        for shape in shapes:
            if shape.shape_type == 6:  # GroupShape
                self.find_and_extract_text(shape.shapes)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell_idx in range(len(row.cells) - 1):
                        key_cell = row.cells[cell_idx]
                        match = find_closest_key(key_cell.text.strip())
                        if match:
                            self.elems[match] = row.cells[cell_idx + 1]

    def generate_ppt(self, content, email_id):
        """Generate PowerPoint presentation with Gemini-extracted info"""

        self.elems = {"Contractor": None, 'Consultant': None, 'LOA Date': None}
        try:
            if os.path.exists(self.config['ppt_template']):
                prs = Presentation(self.config['ppt_template'])
                self.find_and_extract_text(prs.slides[1].shapes)
            else:
                raise FileNotFoundError("PPT Not Found")
            # table.cell(
            #     1, 1).text = content['extracted_info']['contractor'] or "Not identified"
            #
            # table.cell(2, 0).text = "Consultant/Advisor"
            # table.cell(
            #     2, 1).text = content['extracted_info']['consultant'] or "Not identified"
            #
            # table.cell(3, 0).text = "LOA/Approval Date"
            # table.cell(
            #     3, 1).text = content['extracted_info']['loa_date'] or "Not specified"
            self.elems['Contractor'].text = content['extracted_info']['contractor']
            self.elems['Consultant'].text = content['extracted_info']['consultant']
            self.elems['LOA Date'].text = content['extracted_info']['loa_date']
            for k, v in self.elems.items():
                v.text_frame.paragraphs[0].runs[0].font.size = Pt(15)
                # v.text_frame.paragraphs[0].runs[0].font.name = "Ca ta"

            filename = f"processed/analysis_{email_id}_{
                datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            prs.save(filename)
            return filename

        except Exception as error:
            print(f"Error generating PPT: {error}")
            return None

    def mark_as_read(self, service, message_id):
        """Mark email as read"""
        pass
        # try:
        #     service.users().messages().modify(
        #         userId='me',
        #         id=message_id,
        #         body={'removeLabelIds': ['UNREAD']}
        #     ).execute()
        # except Exception as error:
        #     print(f"Error marking email as read: {error}")


# Initialize processor
processor = GmailPPTProcessor()


def credentials_to_dict(credentials):
    """Convert credentials to dictionary"""
    return {
        'token': credentials.token,
        'refresh_token': credentials.refresh_token,
        'token_uri': credentials.token_uri,
        'client_id': credentials.client_id,
        'client_secret': credentials.client_secret,
        'scopes': credentials.scopes
    }


def dict_to_credentials(creds_dict):
    """Convert dictionary to credentials"""
    return Credentials(**creds_dict)


@app.route('/')
def index():
    """Main page"""
    if 'credentials' not in session:
        return render_template('index.html', logged_in=False)
    else:
        return render_template('index.html', logged_in=True)


@app.route('/login')
def login():
    """Start OAuth flow"""
    flow = Flow.from_client_secrets_file(
        CLIENT_SECRETS_FILE,
        scopes=SCOPES,
        redirect_uri=url_for('oauth2callback', _external=True)
    )

    authorization_url, state = flow.authorization_url(
        access_type='offline',
        include_granted_scopes='true',
        prompt='consent'
    )

    session['state'] = state
    return redirect(authorization_url)


@app.route('/oauth2callback')
def oauth2callback():
    """Handle OAuth callback"""
    state = session['state']

    flow = Flow.from_client_secrets_file(
        CLIENT_SECRETS_FILE,
        scopes=SCOPES,
        state=state,
        redirect_uri=url_for('oauth2callback', _external=True)
    )

    flow.fetch_token(authorization_response=request.url)

    credentials = flow.credentials
    session['credentials'] = credentials_to_dict(credentials)

    flash('Successfully logged in to Gmail!', 'success')
    return redirect(url_for('index'))


@app.route('/logout')
def logout():
    """Logout user"""
    session.clear()
    flash('Logged out successfully!', 'info')
    return redirect(url_for('index'))


@app.route('/process_emails')
def process_emails():
    """Process emails and generate PPTs"""
    if 'credentials' not in session:
        flash('Please login first!', 'error')
        return redirect(url_for('index'))

    try:
        credentials = dict_to_credentials(session['credentials'])
        service = processor.get_gmail_service(credentials)

        messages = processor.get_unread_emails(service)
        processed_count = 0
        total_emails = len(messages)

        results = []
        print(messages)

        for message in messages[:3]:  # Limit to 10 emails per run
            message_id = message['id']

            if message_id in processor.processed_emails:
                continue

            email_content = processor.get_email_content(service, message_id)
            if not email_content:
                continue

            result = {
                'id': message_id,
                'subject': email_content['subject'],
                'sender': email_content['sender'],
                'processed': False,
                'ppt_file': None
            }
            print(result)

            if processor.check_for_trigger_text(email_content):
                ppt_content = processor.extract_ppt_content(email_content)
                ppt_file = processor.generate_ppt(ppt_content, message_id)

                if ppt_file:
                    processed_count += 1
                    result['processed'] = True
                    result['ppt_file'] = os.path.basename(ppt_file)

            processor.processed_emails.add(message_id)
            processor.mark_as_read(service, message_id)
            results.append(result)

        processor.save_processed_emails()

        return render_template('process_results.html',
                               results=results,
                               processed_count=processed_count,
                               total_emails=total_emails)

    except Exception as e:
        flash(f'Error processing emails: {str(e)}', 'error')
        return redirect(url_for('index'))


@app.route('/ppt_files')
def ppt_files():
    """List all PPT files"""
    try:
        files = []
        processed_dir = 'processed'

        if os.path.exists(processed_dir):
            for filename in os.listdir(processed_dir):
                if filename.endswith('.pptx'):
                    filepath = os.path.join(processed_dir, filename)
                    file_info = {
                        'name': filename,
                        'size': os.path.getsize(filepath),
                        'created': datetime.fromtimestamp(os.path.getctime(filepath)).strftime('%Y-%m-%d %H:%M:%S')
                    }
                    files.append(file_info)

        files.sort(key=lambda x: x['created'], reverse=True)
        return render_template('ppt_files.html', files=files)

    except Exception as e:
        flash(f'Error listing files: {str(e)}', 'error')
        return redirect(url_for('index'))


@app.route('/download/<filename>')
def download_file(filename):
    """Download PPT file"""
    try:
        filepath = os.path.join('processed', filename)
        if os.path.exists(filepath) and filename.endswith('.pptx'):
            return send_file(filepath, as_attachment=True)
        else:
            flash('File not found!', 'error')
            return redirect(url_for('ppt_files'))
    except Exception as e:
        flash(f'Error downloading file: {str(e)}', 'error')
        return redirect(url_for('ppt_files'))


def get_font_data():
    # Path to your font directory
    font_dir = os.path.join(
        app.static_folder, 'font/Proxima Nova Complete Collection')

    # Find all font files
    font_files = []
    for ext in ['ttf', 'otf', 'woff', 'woff2']:
        font_files.extend(glob.glob(os.path.join(font_dir, f'*.{ext}')))

    # Extract just the filenames
    font_files = [os.path.basename(f) for f in font_files]

    return font_files


@app.context_processor
def inject_font_data():
    return {'font_files': get_font_data()}


if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000, ssl_context='adhoc')
